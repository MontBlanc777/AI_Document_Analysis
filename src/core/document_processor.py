import os
import base64
import re
import requests
from typing import Dict, Any, List, Optional, Tuple, BinaryIO, Union
from datetime import datetime
import mimetypes
import json
import PyPDF2
from urllib.parse import urlparse
from src.utils.logging_utils import get_app_logger

# Initialize logger
logger = get_app_logger()

# Try to import various document processing libraries
# These will be used if available, otherwise fallback methods will be used
try:
    from pdf2image import convert_from_path, convert_from_bytes
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False
    logger.warning("pdf2image not available, PDF processing will be limited")

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    import email
    import email.policy
    from email.parser import BytesParser
    EMAIL_AVAILABLE = True
except ImportError:
    EMAIL_AVAILABLE = False

class DocumentProcessor:
    """Class for processing various document types"""
    
    def __init__(self, upload_folder: str):
        """Initialize the document processor"""
        self.upload_folder = upload_folder
        os.makedirs(upload_folder, exist_ok=True)
        logger.info(f"DocumentProcessor initialized with upload folder: {upload_folder}")
    
    def save_uploaded_file(self, file_content: bytes, filename: str) -> str:
        """Save an uploaded file to disk and return the file path"""
        # Check if filename is valid
        if not filename or filename.strip() == '':
            logger.error("Attempted to save file with empty filename")
            raise ValueError("Filename cannot be empty")
            
        # Log the upload attempt
        logger.info(f"Saving file: {filename}, content size: {len(file_content)} bytes")
        
        # Generate a unique filename
        timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
        unique_filename = f"{timestamp}_{filename}"
        file_path = os.path.join(self.upload_folder, unique_filename)
        
        # Ensure upload directory exists
        os.makedirs(self.upload_folder, exist_ok=True)
        
        # Log the file path
        logger.info(f"Saving to path: {file_path}")
        
        # Save the file
        try:
            with open(file_path, "wb") as f:
                f.write(file_content)
            
            # Verify file was saved
            if os.path.exists(file_path):
                file_size = os.path.getsize(file_path)
                logger.info(f"File saved successfully: {file_path}, size: {file_size} bytes")
            else:
                logger.error(f"File was not found after saving: {file_path}")
                raise FileNotFoundError(f"File was not found after saving: {file_path}")
                
            return file_path
        except Exception as e:
            logger.error(f"Error saving file {filename}: {str(e)}")
            raise
        
    def is_valid_url(self, url: str) -> bool:
        """Check if a URL is valid"""
        try:
            result = urlparse(url)
            return all([result.scheme, result.netloc])
        except:
            return False
            
    def download_from_url(self, url: str) -> Tuple[str, bytes, str]:
        """Download content from a URL and return filename, content and mime type"""
        try:
            response = requests.get(url, timeout=30)
            response.raise_for_status()
            
            # Try to get filename from URL or Content-Disposition header
            filename = None
            content_disposition = response.headers.get('Content-Disposition')
            if content_disposition:
                filename_match = re.search(r'filename="?([^"]+)"?', content_disposition)
                if filename_match:
                    filename = filename_match.group(1)
            
            if not filename:
                # Extract filename from URL
                parsed_url = urlparse(url)
                path = parsed_url.path
                filename = os.path.basename(path)
                
            # If still no filename, use the domain with timestamp
            if not filename or filename == '':
                domain = urlparse(url).netloc
                timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
                filename = f"{domain}_{timestamp}.html"
            
            # Get content type
            content_type = response.headers.get('Content-Type', 'text/html')
            
            return filename, response.content, content_type
        except Exception as e:
            raise Exception(f"Failed to download from URL: {str(e)}")
            
    def process_url(self, url: str) -> Tuple[str, str, str]:
        """Process a URL by downloading its content and saving it"""
        filename, content, mime_type = self.download_from_url(url)
        file_path = self.save_uploaded_file(content, filename)
        return file_path, filename, mime_type
    
    def get_mime_type(self, filename: str) -> str:
        """Get the MIME type of a file based on its extension"""
        mime_type, _ = mimetypes.guess_type(filename)
        if mime_type is None:
            # Default to binary if we can't determine the type
            mime_type = "application/octet-stream"
        return mime_type
    
    def process_document(self, file_path: str) -> Dict[str, Any]:
        """Process a document based on its MIME type"""
        if not os.path.exists(file_path):
            return {"error": f"File not found: {file_path}"}
        
        mime_type = self.get_mime_type(file_path)
        
        # Process based on MIME type
        if mime_type == "application/pdf":
            return self.process_pdf(file_path)
        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return self.process_docx(file_path)
        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return self.process_pptx(file_path)
        elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return self.process_xlsx(file_path)
        elif mime_type in ["message/rfc822", "application/vnd.ms-outlook"]:
            return self.process_email(file_path)
        elif mime_type.startswith("image/"):
            return self.process_image(file_path)
        elif mime_type.startswith("text/"):
            return self.process_text(file_path)
        else:
            # For unknown types, just return basic file info
            return self.get_file_info(file_path)
    
    def get_file_info(self, file_path: str) -> Dict[str, Any]:
        """Get basic file information"""
        try:
            file_size = os.path.getsize(file_path)
            file_name = os.path.basename(file_path)
            mime_type = self.get_mime_type(file_path)
            created_time = os.path.getctime(file_path)
            modified_time = os.path.getmtime(file_path)
            
            return {
                "file_path": file_path,
                "file_name": file_name,
                "mime_type": mime_type,
                "file_size": file_size,
                "created_time": datetime.fromtimestamp(created_time).isoformat(),
                "modified_time": datetime.fromtimestamp(modified_time).isoformat(),
            }
        except Exception as e:
            return {"error": str(e), "file_path": file_path}
    
    def process_pdf(self, file_path: str) -> Dict[str, Any]:
        """Process a PDF document"""
        try:
            result = self.get_file_info(file_path)
            
            # Extract text from PDF using PyPDF2 if available
            try:
                pdf_text = []
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    num_pages = len(pdf_reader.pages)
                    result["page_count"] = num_pages
                    
                    # Extract text from each page
                    for page_num in range(num_pages):
                        page = pdf_reader.pages[page_num]
                        page_text = page.extract_text()
                        if page_text:
                            pdf_text.append(f"Page {page_num+1}:\n{page_text}")
                
                # Join all pages' text
                if pdf_text:
                    result["text_content"] = "\n\n".join(pdf_text)
                    logger.info(f"Extracted {len(result['text_content'])} characters of text from PDF: {os.path.basename(file_path)}")
            except ImportError:
                logger.warning("PyPDF2 not available, falling back to basic PDF processing")
            except Exception as e:
                logger.error(f"Error extracting text from PDF: {str(e)}")
                result["text_extraction_error"] = str(e)
            
            # Extract images if pdf2image is available
            if PDF2IMAGE_AVAILABLE:
                try:
                    # Convert first few pages to images
                    images = convert_from_path(file_path, first_page=1, last_page=3)
                    
                    # Save images to temporary files
                    image_paths = []
                    for i, image in enumerate(images):
                        image_path = os.path.join(self.upload_folder, f"pdf_page_{i+1}_{os.path.basename(file_path)}.jpg")
                        image.save(image_path, "JPEG")
                        image_paths.append(image_path)
                    
                    result["extracted_images"] = image_paths
                    result["num_images_extracted"] = len(image_paths)
                except Exception as e:
                    result["image_extraction_error"] = str(e)
            
            return result
        except Exception as e:
            return {"error": str(e), "file_path": file_path}
    
    def process_docx(self, file_path: str) -> Dict[str, Any]:
        """Process a Word document"""
        try:
            result = self.get_file_info(file_path)
            
            if DOCX_AVAILABLE:
                try:
                    doc = docx.Document(file_path)
                    
                    # Extract text
                    paragraphs = [p.text for p in doc.paragraphs]
                    result["text_content"] = "\n".join(paragraphs)
                    result["paragraph_count"] = len(paragraphs)
                    
                    # Extract tables
                    tables = []
                    for table in doc.tables:
                        table_data = []
                        for row in table.rows:
                            row_data = [cell.text for cell in row.cells]
                            table_data.append(row_data)
                        tables.append(table_data)
                    
                    result["tables"] = tables
                    result["table_count"] = len(tables)
                except Exception as e:
                    result["docx_extraction_error"] = str(e)
            
            return result
        except Exception as e:
            return {"error": str(e), "file_path": file_path}
    
    def process_pptx(self, file_path: str) -> Dict[str, Any]:
        """Process a PowerPoint presentation"""
        try:
            result = self.get_file_info(file_path)
            
            if PPTX_AVAILABLE:
                try:
                    prs = Presentation(file_path)
                    
                    # Extract slides
                    slides = []
                    for slide in prs.slides:
                        slide_text = ""
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                slide_text += shape.text + "\n"
                        slides.append(slide_text.strip())
                    
                    result["slides"] = slides
                    result["slide_count"] = len(slides)
                except Exception as e:
                    result["pptx_extraction_error"] = str(e)
            
            return result
        except Exception as e:
            return {"error": str(e), "file_path": file_path}
    
    def process_xlsx(self, file_path: str) -> Dict[str, Any]:
        """Process an Excel spreadsheet"""
        try:
            result = self.get_file_info(file_path)
            
            if XLSX_AVAILABLE:
                try:
                    workbook = openpyxl.load_workbook(file_path, read_only=True)
                    
                    # Extract sheets
                    sheets = {}
                    for sheet_name in workbook.sheetnames:
                        sheet = workbook[sheet_name]
                        sheet_data = []
                        for row in sheet.iter_rows(values_only=True, max_row=100):  # Limit to 100 rows
                            sheet_data.append(list(row))
                        sheets[sheet_name] = sheet_data
                    
                    result["sheets"] = sheets
                    result["sheet_count"] = len(sheets)
                except Exception as e:
                    result["xlsx_extraction_error"] = str(e)
            
            return result
        except Exception as e:
            return {"error": str(e), "file_path": file_path}
    
    def process_email(self, file_path: str) -> Dict[str, Any]:
        """Process an email message"""
        try:
            result = self.get_file_info(file_path)
            
            if EMAIL_AVAILABLE:
                try:
                    with open(file_path, "rb") as f:
                        msg = BytesParser(policy=email.policy.default).parse(f)
                    
                    # Extract headers
                    headers = {
                        "From": msg.get("From", ""),
                        "To": msg.get("To", ""),
                        "Subject": msg.get("Subject", ""),
                        "Date": msg.get("Date", ""),
                    }
                    result["headers"] = headers
                    
                    # Extract body
                    body = ""
                    if msg.is_multipart():
                        for part in msg.iter_parts():
                            if part.get_content_type() == "text/plain":
                                body += part.get_content()
                    else:
                        body = msg.get_content()
                    
                    result["body"] = body
                    
                    # Extract attachments
                    attachments = []
                    if msg.is_multipart():
                        for part in msg.iter_parts():
                            if part.get_content_disposition() == "attachment":
                                filename = part.get_filename()
                                if filename:
                                    attachment_path = os.path.join(self.upload_folder, f"attachment_{filename}")
                                    with open(attachment_path, "wb") as f:
                                        f.write(part.get_payload(decode=True))
                                    attachments.append({
                                        "filename": filename,
                                        "path": attachment_path,
                                        "mime_type": part.get_content_type(),
                                    })
                    
                    result["attachments"] = attachments
                    result["attachment_count"] = len(attachments)
                except Exception as e:
                    result["email_extraction_error"] = str(e)
            
            return result
        except Exception as e:
            return {"error": str(e), "file_path": file_path}
    
    def process_image(self, file_path: str) -> Dict[str, Any]:
        """Process an image file"""
        try:
            result = self.get_file_info(file_path)
            
            # In a real implementation, you would extract image metadata and possibly OCR text
            # For this example, we'll just return the basic info
            
            return result
        except Exception as e:
            return {"error": str(e), "file_path": file_path}
    
    def process_text(self, file_path: str) -> Dict[str, Any]:
        """Process a text file"""
        try:
            result = self.get_file_info(file_path)
            
            # Extract text content
            try:
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f.read()
                result["text_content"] = content
                result["line_count"] = content.count("\n") + 1
            except UnicodeDecodeError:
                # Try with a different encoding
                with open(file_path, "r", encoding="latin-1") as f:
                    content = f.read()
                result["text_content"] = content
                result["line_count"] = content.count("\n") + 1
            
            return result
        except Exception as e:
            return {"error": str(e), "file_path": file_path}
    
    def encode_file_for_gemini(self, file_path: str) -> Tuple[str, str]:
        """Encode a file for use with Gemini API"""
        mime_type = self.get_mime_type(file_path)
        
        with open(file_path, "rb") as f:
            file_content = f.read()
        
        # For text files, try to decode as text
        if mime_type.startswith("text/"):
            try:
                text_content = file_content.decode("utf-8")
                return mime_type, text_content
            except UnicodeDecodeError:
                # If decoding fails, fall back to base64
                pass
        
        # For binary files, encode as base64
        encoded_content = base64.b64encode(file_content).decode("utf-8")
        return mime_type, encoded_content